#!/usr/bin/env python3
"""Build presentation/deck.pptx — the PowerPoint twin of deck.html.

Same visual identity (#CC0000 system, light theme), same 29-slide structure.
Content slides reveal component-by-component via click-entrance animations
(PowerPoint timing XML, injected — python-pptx has no animation API).
The walkthrough (slides 14-23) advances state per slide, exactly like the
HTML deck; the C4 vote buttons are real hyperlinks to the YES / NO slides.

Fonts: Noto Sans for text (falls back gracefully if not installed),
Consolas for code/facts/traces (ships with Office).
Regenerate:  python3 presentation/build_pptx.py
"""
import copy, re, pathlib
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn
from lxml import etree

# ── identity tokens ────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0xCC, 0x00, 0x00)
PRIMARY_D = RGBColor(0x99, 0x00, 0x00)
PRIMARY_T = RGBColor(0xFF, 0xEA, 0xEA)
SUCCESS   = RGBColor(0x2D, 0x6A, 0x4F)
SUCCESS_T = RGBColor(0xF0, 0xFA, 0xF5)
WARN_TX   = RGBColor(0xB5, 0x4A, 0x00)
WARN_T    = RGBColor(0xFF, 0xF3, 0xEB)
INK       = RGBColor(0x1A, 0x1A, 0x1A)
INK2      = RGBColor(0x4A, 0x4A, 0x4A)
INK3      = RGBColor(0x6B, 0x6B, 0x6B)
BORDER    = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_EDGE = RGBColor(0xC7, 0xC7, 0xC7)
SURFACE   = RGBColor(0xF5, 0xF5, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DIM_RED   = RGBColor(0xE5, 0xA3, 0xA3)   # "dim" variants (HTML used opacity)
DIM_GREEN = RGBColor(0xA9, 0xC9, 0xBA)
DIM_WARN  = RGBColor(0xE3, 0xC2, 0x9E)
PAPER_BG  = RGBColor(0xFB, 0xF8, 0xF1)
PAPER_BR  = RGBColor(0xB9, 0xAE, 0x8F)
PAPER_TX  = RGBColor(0x5C, 0x53, 0x40)

SANS, MONO = "Noto Sans", "Consolas"
PX = 9525                                # 1 CSS px (96dpi) in EMU

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(1280*PX), Emu(720*PX)
BLANK = prs.slide_layouts[6]

def E(v): return Emu(round(v*PX))

# ── low-level helpers ──────────────────────────────────────────────────────
def _style_runs(para, runs, size, align=None):
    if align: para.alignment = align
    for txt, kw in runs:
        r = para.add_run(); r.text = txt
        f = r.font
        f.name = kw.get("font", SANS); f.size = Pt(kw.get("size", size))
        f.bold = kw.get("bold", False); f.italic = kw.get("italic", False)
        f.color.rgb = kw.get("color", INK)
        f.language_id = MSO_LANGUAGE_ID.ENGLISH_US
        if kw.get("strike"): r.font._rPr.set("strike", "sngStrike")

def textbox(slide, x, y, w, h, paras, size=13, align=None, anchor=MSO_ANCHOR.TOP,
            wrap=True, spacing=None):
    """paras: list of para-specs; each is list of (text, kwargs) runs."""
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if spacing is not None: p.line_spacing = spacing
        _style_runs(p, runs, size)
    return tb

def box(slide, x, y, w, h, fill=WHITE, line=BORDER, lw=1.5, radius=.12,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None, shadow=False):
    s = slide.shapes.add_shape(shape, E(x), E(y), E(w), E(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    if dash:
        ln = s.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"}); ln.append(d)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    tf = s.text_frame
    tf.margin_left = tf.margin_right = E(8)
    tf.margin_top = tf.margin_bottom = E(4)
    return s

def set_text(shape, paras, size=12, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame; tf.vertical_anchor = anchor
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _style_runs(p, runs, size, align)
    return shape

def chip(slide, x, y, w, h, runs, fill=WHITE, line=BORDER, txtsize=11):
    c = box(slide, x, y, w, h, fill=fill, line=line, lw=1.25, radius=.5)
    set_text(c, [runs], size=txtsize)
    return c

def line_seg(slide, x1, y1, x2, y2, color, w=2.0, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    c.shadow.inherit = False
    if arrow:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                  {"type": "triangle", "w": "med", "len": "med"}))
    return c

def path_segs(d):
    """parse the tiny M/V/H/L path syntax used by deck.html edge data."""
    toks = re.findall(r"([MVHL])\s*([\d.,\s]+)", d)
    segs, cx, cy = [], 0, 0
    for cmd, val in toks:
        nums = [float(n) for n in re.split(r"[,\s]+", val.strip()) if n]
        if cmd == "M": cx, cy = nums
        elif cmd == "V": segs.append((cx, cy, cx, nums[0])); cy = nums[0]
        elif cmd == "H": segs.append((cx, cy, nums[0], cy)); cx = nums[0]
        elif cmd == "L": segs.append((cx, cy, nums[0], nums[1])); cx, cy = nums
    return segs

def edge(slide, d, color, w=2.0):
    shapes, segs = [], path_segs(d)
    for i, (x1, y1, x2, y2) in enumerate(segs):
        shapes.append(line_seg(slide, x1, y1, x2, y2, color, w, arrow=(i == len(segs)-1)))
    return shapes

def slide_no(slide, n):
    textbox(slide, 1180, 694, 80, 20,
            [[(f"{n} / 28", {"size": 9, "color": INK3})]], align=PP_ALIGN.RIGHT)

def header(slide, kicker, title, tsize=25):
    shapes = []
    shapes += [textbox(slide, 44, 30, 900, 20,
        [[(kicker.upper(), {"size": 10.5, "bold": True, "color": PRIMARY})]])]
    shapes += [textbox(slide, 44, 52, 1000, 46,
        [[(title, {"size": tsize, "bold": True})]])]
    wm = box(slide, 1130, 34, 106, 30, fill=PRIMARY, line=None, radius=.28)
    set_text(wm, [[("ResQRules", {"size": 12.5, "bold": True, "color": WHITE})]])
    shapes.append(wm)
    return shapes

def bullet(slide, x, y, w, runs, size=17):
    sq = box(slide, x, y+7, 8, 8, fill=PRIMARY, line=None, radius=.3)
    tb = textbox(slide, x+20, y, w-20, 60, [runs], size=size, spacing=1.1)
    return [sq, tb]

def statement(slide, x, y, w, text, size=19, color=PRIMARY_D):
    bar = box(slide, x, y, 5, 46, fill=PRIMARY, line=None, radius=.3)
    tb = textbox(slide, x+16, y+2, w-16, 60,
                 [[(text, {"size": size, "bold": True, "color": color})]], spacing=1.1)
    return [bar, tb]

def codebox(slide, x, y, w, h, paras, size=10.5, font=MONO):
    b = box(slide, x, y, w, h, fill=SURFACE, line=BORDER, lw=1.25, radius=.06)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = E(12); tf.margin_top = tf.margin_bottom = E(8)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        _style_runs(p, [(t, {**kw, "font": kw.get("font", font)}) for t, kw in runs], size)
    return b

def red_slide():
    s = prs.slides.add_slide(BLANK)
    bg = box(s, -2, -2, 1284, 724, fill=PRIMARY, line=None, radius=0,
             shape=MSO_SHAPE.RECTANGLE)
    return s

# ── click-appear animation (timing XML injection) ──────────────────────────
def animate(slide, groups):
    """groups: list of clicks; each click = list of shapes that appear together."""
    groups = [[sp for sp in g if sp is not None] for g in groups]
    groups = [g for g in groups if g]
    if not groups: return
    nid = [1]
    def i_(): nid[0] += 1; return nid[0]
    P = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" ' \
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    clicks = []
    for g in groups:
        effects = []
        for j, sp in enumerate(g):
            node_t = "clickEffect" if j == 0 else "withEffect"
            effects.append(f"""
              <p:par><p:cTn id="{i_()}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="{node_t}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst><p:set><p:cBhvr>
                  <p:cTn id="{i_()}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                  <p:tgtEl><p:spTgt spid="{sp.shape_id}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
              </p:cTn></p:par>""")
        clicks.append(f"""
          <p:par><p:cTn id="{i_()}" fill="hold">
            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="{i_()}" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>{''.join(effects)}</p:childTnLst>
            </p:cTn></p:par></p:childTnLst>
          </p:cTn></p:par>""")
    xml = f"""<p:timing {P}><p:tnLst><p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        <p:seq concurrent="1" nextAc="seek">
          <p:cTn id="{i_()}" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
            {''.join(clicks)}
          </p:childTnLst></p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
        </p:seq>
      </p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""
    # fix duplicated id=1 on root: renumber root to a fresh unique id? PowerPoint
    # requires tmRoot id unique within tree; ids above start at 2 via i_(), root=1. OK.
    el = etree.fromstring(xml.encode())
    slide._element.append(el)

LINKS = []   # (shape, slide_key) resolved after all slides exist
SLIDES = {}  # key -> slide

# ═══════════════════════════ SLIDES ════════════════════════════════════════

# s01 — title (red)
s = red_slide(); SLIDES["s01"] = s
textbox(s, 140, 200, 1000, 110, [[("ResQRules", {"size": 66, "bold": True, "color": WHITE})]],
        align=PP_ALIGN.CENTER)
textbox(s, 190, 330, 900, 90, [[(
    "A knowledge-based system that walks a first responder through Red Crescent "
    "emergency protocols — and reacts when the patient doesn't follow the flowchart.",
    {"size": 19, "color": RGBColor(0xFF, 0xE2, 0xE2)})]], align=PP_ALIGN.CENTER, spacing=1.25)
textbox(s, 190, 440, 900, 30, [[(
    "Knowledge-Based Systems · course project · built on the Experta forward-chaining engine",
    {"size": 13, "color": RGBColor(0xFF, 0xB9, 0xB9)})]], align=PP_ALIGN.CENTER)
textbox(s, 190, 500, 900, 30, [[(
    "M.GH · M.SH · M.JA · M.Y.D · M.Marwan  —  supervised by Miss Omama",
    {"size": 13, "color": RGBColor(0xFF, 0xB9, 0xB9)})]], align=PP_ALIGN.CENTER)
slide_no(s, 1)

# s02 — the problem
s = prs.slides.add_slide(BLANK); SLIDES["s02"] = s
header(s, "The problem", "The patient doesn't follow the flowchart")
g_paper = []
g_paper.append(box(s, 60, 150, 350, 440, fill=PAPER_BG, line=PAPER_BR, lw=1.2, radius=.03))
g_paper.append(textbox(s, 78, 164, 320, 18, [[(
    "SARC SOP · CATASTROPHIC HEMORRHAGE · CHART #10",
    {"size": 8.5, "bold": True, "color": RGBColor(0x8A, 0x7F, 0x63)})]]))
pb = [("Control of external bleeding: direct compression…", 0),
      ("Active bleeding?", 1), ("Second compression bandage on top of the first", 0),
      ("Active bleeding?", 1), ("Installation of tourniquet…", 0),
      ("Urgent and safe transport…", 0)]
yy = 192
for txt, dashed in pb:
    bx = box(s, 95, yy, 280, 34, fill=WHITE, line=PAPER_BR, lw=1.1, radius=.06,
             dash=("dash" if dashed else None))
    set_text(bx, [[(txt, {"size": 9.5, "color": PAPER_TX})]])
    g_paper.append(bx)
    if txt != pb[-1][0]:
        g_paper.append(textbox(s, 220, yy+34, 40, 16, [[("↓", {"size": 10, "color": PAPER_BR})]]))
    yy += 62
here = box(s, 330, 300, 170, 52, fill=PRIMARY, line=None, radius=.15)
set_text(here, [[("you are here…", {"size": 10.5, "bold": True, "color": WHITE})],
                [("but he just stopped breathing", {"size": 10.5, "bold": True, "color": WHITE})]])
here.rotation = 3
b1 = bullet(s, 540, 200, 690, [("The Red Crescent SOP is ", {}), ("12 paper flowcharts", {"bold": True}),
                               (" — expert, validated, and static.", {})])
b2 = bullet(s, 540, 270, 690, [("Under adrenaline, the responder must find the right chart, the right box, "
                                "and follow the arrows correctly.", {})])
b3 = bullet(s, 540, 350, 690, [("Meanwhile the patient ", {}), ("deteriorates mid-protocol", {"bold": True}),
                               (" — a new life threat appears while you're inside another box.", {})])
st = statement(s, 540, 450, 690, "Paper can't re-route. That gap is what we built.")
slide_no(s, 2)
animate(s, [g_paper, [here], b1, b2, b3, st])

# s03 — what it is
s = prs.slides.add_slide(BLANK); SLIDES["s03"] = s
header(s, "ResQRules", "Protocols as software that reacts")
steps = [("Experta engine", "terminal Q&A · the reasoning"),
         ("FastAPI", "stable REST contract"),
         ("Flutter app", "bilingual EN/AR · in the field")]
g_steps = []
x = 210
for i, (t, d) in enumerate(steps):
    bx = box(s, x, 250, 250, 92, fill=WHITE, line=INK2, lw=1.75, radius=.14)
    set_text(bx, [[(t, {"size": 16, "bold": True})], [(d, {"size": 11, "color": INK3})]])
    grp = [bx]
    if i < 2:
        grp.append(textbox(s, x+254, 278, 40, 36, [[("→", {"size": 20, "bold": True, "color": INK3})]]))
    g_steps.append(grp)
    x += 300
chips_g = []
cx = 150
for t, red in [("Field responders — guided, interruptible protocols", 0),
               ("Trainees — safe practice", 0), ("This room — a working KBS case study", 1)]:
    w = 24 + len(t)*7
    chips_g.append(chip(s, cx, 420, w, 34,
                        [(t, {"size": 11.5, "bold": True, "color": PRIMARY_D if red else INK2})],
                        fill=PRIMARY_T if red else WHITE, line=PRIMARY if red else BORDER))
    cx += w + 18
cite = chip(s, 44, 660, 640, 30, [("5 protocols in scope: #1 assessment · #3 CPR adult · #7 choking · #8 airway · #10 hemorrhage",
             {"size": 10.5, "color": INK2})], fill=SURFACE)
slide_no(s, 3)
animate(s, g_steps + [chips_g, [cite]])

# s04 — phases (lanes appear one by one, then the exception)
s = prs.slides.add_slide(BLANK); SLIDES["s04"] = s
header(s, "Domain — how a paramedic thinks", "Four assessment phases — and one structural exception")
lanes = [("1", "Preliminary — SSS", "Safety · Scene · Situation"),
         ("2", "Primary", "AVPU + ABCDE: Airway · Breathing · Circulation · Disability · Expose"),
         ("3", "Secondary", "SAMPLER · OPQRST · GCS · PERRLA · vitals"),
         ("4", "Continuous", "Reassess & orient to the right care center")]
g_lanes = []
x = 44
for n, t, d in lanes:
    lane = box(s, x, 160, 288, 150, fill=WHITE, line=BORDER, lw=1.75, radius=.09)
    num = box(s, x+16, 176, 30, 30, fill=INK2, line=None, shape=MSO_SHAPE.OVAL, radius=None)
    set_text(num, [[(n, {"size": 13, "bold": True, "color": WHITE})]])
    tt = textbox(s, x+16, 214, 258, 26, [[(t, {"size": 15, "bold": True})]])
    dd = textbox(s, x+16, 244, 258, 58, [[(d, {"size": 11.5, "color": INK2})]], spacing=1.15)
    g_lanes.append([lane, num, tt, dd])
    x += 300
exc = []
exc.append(box(s, 44, 360, 1192, 120, fill=PRIMARY_T, line=PRIMARY, lw=2, radius=.1))
exc.append(textbox(s, 70, 384, 240, 70, [[("Life-threatening", {"size": 16, "bold": True, "color": PRIMARY_D})],
                                         [("bleeding?  (gen_04)", {"size": 16, "bold": True, "color": PRIMARY_D})]]))
exc.append(textbox(s, 330, 378, 880, 90, [
    [("Asked ", {"size": 12.5}), ("between phases 1 and 2 — before ABCDE ever runs", {"size": 12.5, "bold": True}),
     (" — jumping straight to the Catastrophic Hemorrhage protocol. The chart encodes clinical precedence ", {"size": 12.5}),
     ("positionally", {"size": 12.5, "bold": True}), (".", {"size": 12.5})],
    [("Our engine encodes the same precedence as salience", {"size": 12.5, "bold": True}),
     (" — coming up in Inference I.", {"size": 12.5})]], spacing=1.2))
cite = chip(s, 44, 660, 560, 30, [("SARC chart #1 · p.1 · gen_04 → jump to hemorrhage (drawn on the chart)",
            {"size": 10.5, "color": INK2})], fill=SURFACE)
slide_no(s, 4)
animate(s, g_lanes + [exc, [cite]])

# s05 — sources
s = prs.slides.add_slide(BLANK); SLIDES["s05"] = s
header(s, "Knowledge sources", "Expert-validated knowledge — never our invention")
g1 = [box(s, 44, 170, 580, 260, fill=WHITE, line=BORDER, lw=1.5, radius=.06)]
g1.append(textbox(s, 70, 194, 520, 30, [[("SARC official SOP flowcharts", {"size": 17, "bold": True})]]))
g1.append(textbox(s, 70, 230, 520, 26, [[("The 12-chart first-responder deck; 5 in scope for v1:", {"size": 12.5, "color": INK2})]]))
cy = 268
for t in ["#1 General Assessment", "#3 CPR Adult", "#7 Heimlich / Choking", "#8 Upper Airway", "#10 Catastrophic Hemorrhage"]:
    g1.append(chip(s, 70, cy, 250, 28, [(t, {"size": 11, "bold": True, "color": PRIMARY_D})],
                   fill=WHITE, line=PRIMARY)); cy += 34
g2 = [box(s, 656, 170, 580, 260, fill=WHITE, line=BORDER, lw=1.5, radius=.06)]
g2.append(textbox(s, 682, 194, 520, 30, [[("An expert in the loop", {"size": 17, "bold": True})]]))
g2.append(textbox(s, 682, 232, 520, 170, [[
    ("M.GH — paramedic on the team. ", {"size": 13, "bold": True}),
    ("Collected the source charts, resolved ambiguities during extraction, and validated the "
     "transcribed knowledge. Knowledge acquisition with the expert at the table, not by email.",
     {"size": 13, "color": INK2})]], spacing=1.3))
st = statement(s, 44, 480, 900, "Rule zero: transcribe the chart — never “improve” medical content.")
slide_no(s, 5)
animate(s, [g1, g2, st])

# s06 — rejected design
s = prs.slides.add_slide(BLANK); SLIDES["s06"] = s
header(s, "Design — what we rejected", "One @Rule per box is not a knowledge base")
wall = []
ids = ["gen_01","gen_02","gen_03","gen_04","chk_01","chk_02","chk_03","chk_04",
       "cpr_01","cpr_02","cpr_03","cpr_04","hem_01","hem_02","air_01","air_02"]
for i, nid_ in enumerate(ids):
    xx, yy = 50 + (i % 4)*150, 190 + (i // 4)*44
    c = box(s, xx, yy, 140, 32, fill=SURFACE, line=BORDER, lw=1, radius=.12)
    set_text(c, [[(f"@Rule(F(nid='{nid_}'))…", {"size": 8.5, "font": MONO, "color": INK3})]])
    wall.append(c)
wall.append(textbox(s, 50, 372, 590, 34, [[("× 59 rules", {"size": 18, "bold": True, "color": INK3})]],
                    align=PP_ALIGN.CENTER))
xg = [line_seg(s, 60, 195, 630, 395, PRIMARY, 5), line_seg(s, 630, 195, 60, 395, PRIMARY, 5)]
b1 = bullet(s, 700, 190, 540, [("Knowledge ", {}), ("trapped inside code", {"bold": True}),
                               (" — every medical revision is a code change.", {})], size=15)
b2 = bullet(s, 700, 255, 540, [("The rule engine adds nothing: ", {}), ("no inference is demonstrated", {"bold": True}),
                               (".", {})], size=15)
b3 = bullet(s, 700, 320, 540, [("Experta reduced to a ", {}), ("flowchart costume", {"bold": True}), (".", {})], size=15)
st = statement(s, 700, 400, 540, "“Exactly like a web app with static hardcoded pages — and no database.”", size=16)
slide_no(s, 6)
animate(s, [wall, xg, b1, b2, b3, st])

# s07 — built design A
s = prs.slides.add_slide(BLANK); SLIDES["s07"] = s
header(s, "Design — what we built", "Knowledge is data; inference is generic")
kb = [box(s, 150, 200, 420, 190, fill=PRIMARY_T, line=PRIMARY, lw=2, radius=.08)]
kb.append(textbox(s, 170, 218, 380, 28, [[("KNOWLEDGE — data/*.json", {"size": 15.5, "bold": True, "color": PRIMARY_D})]], align=PP_ALIGN.CENTER))
kb.append(textbox(s, 170, 250, 380, 24, [[("one tree per chart · 59 nodes · full provenance", {"size": 11, "color": INK3})]], align=PP_ALIGN.CENTER))
cx, cyy = 178, 286
for t in ["general_assessment", "cpr_adult", "choking", "airway", "hemorrhage"]:
    w = 30 + len(t)*6.6
    kb.append(chip(s, cx, cyy, w, 27, [(t, {"size": 10, "font": MONO, "color": PRIMARY_D})], line=PRIMARY))
    cx += w + 10
    if cx > 470: cx, cyy = 178, 322
mid = [textbox(s, 590, 260, 110, 70, [[("⇄", {"size": 24, "bold": True, "color": INK2})],
       [("loads any conforming tree", {"size": 9.5, "color": INK3})]], align=PP_ALIGN.CENTER)]
eng = [box(s, 720, 200, 380, 190, fill=WHITE, line=INK2, lw=2, radius=.08)]
eng.append(textbox(s, 740, 224, 340, 28, [[("INFERENCE — engine.py", {"size": 15.5, "bold": True})]], align=PP_ALIGN.CENTER))
eng.append(textbox(s, 740, 258, 340, 100, [
    [("one KnowledgeEngine", {"size": 12, "color": INK2})],
    [("8 generic rules", {"size": 22, "bold": True, "color": PRIMARY})],
    [("zero chart-specific code", {"size": 12, "color": INK2})]], align=PP_ALIGN.CENTER, spacing=1.25))
contract = [box(s, 150, 420, 950, 54, fill=WHITE, line=BORDER, lw=1.5, radius=.1, dash="dash")]
contract.append(set_text(contract[0], [[("SCHEMA.md — the contract between them   ·   validate.py — fails loudly on any broken edge, orphan node, or bad leaf",
             {"size": 12, "color": INK2})]]))
close = textbox(s, 150, 505, 950, 34, [[("This separation is the textbook definition of a knowledge-based system.",
        {"size": 16.5, "bold": True, "color": PRIMARY_D})]], align=PP_ALIGN.CENTER)
slide_no(s, 7)
animate(s, [kb, mid + eng, contract, [close]])

# s08 — built design B
s = prs.slides.add_slide(BLANK); SLIDES["s08"] = s
header(s, "Design — what we built", "…and the whole product stands on that separation")
top = []
for i, (t, d) in enumerate([("Flutter app", "renders whatever node the engine serves"),
                            ("FastAPI", "contract driven by node type — not by rules")]):
    bx = box(s, 400 + i*250, 150, 230, 84, fill=WHITE, line=INK2, lw=1.75, radius=.1)
    set_text(bx, [[(t, {"size": 14.5, "bold": True})], [(d, {"size": 10, "color": INK3})]])
    top.append(bx)
mid = [textbox(s, 400, 244, 480, 26, [[("▲ sits on top of ▲", {"size": 13, "color": INK3})]], align=PP_ALIGN.CENTER)]
for i, (t, d, k) in enumerate([("KNOWLEDGE (JSON)", "5 charts · 59 nodes", 1),
                               ("INFERENCE (Experta)", "8 generic rules + 3 inference layers", 0)]):
    bx = box(s, 280 + i*370, 280, 350, 80, fill=PRIMARY_T if k else WHITE,
             line=PRIMARY if k else INK2, lw=2, radius=.1)
    set_text(bx, [[(t, {"size": 14.5, "bold": True, "color": PRIMARY_D if k else INK})],
                  [(d, {"size": 10.5, "color": INK3})]])
    mid.append(bx)
stats = []
for i, (n, lab) in enumerate([("8", "generic rules"), ("59", "knowledge nodes"), ("5", "protocols")]):
    stats.append(textbox(s, 400 + i*170, 390, 150, 80,
        [[(n, {"size": 34, "bold": True, "color": PRIMARY})], [(lab, {"size": 12, "color": INK2})]],
        align=PP_ALIGN.CENTER))
claim = chip(s, 250, 500, 780, 40, [("Falsifiable claim: chart #6 = drop one JSON file in data/ — zero code change. Test us in Q&A.",
        {"size": 13, "bold": True, "color": PRIMARY_D})], fill=PRIMARY_T, line=PRIMARY, txtsize=13)
slide_no(s, 8)
animate(s, [top, mid, stats, [claim]])

# s09 — pdf→json
s = prs.slides.add_slide(BLANK); SLIDES["s09"] = s
header(s, "Knowledge acquisition", "Faithfulness is engineered, not promised")
gpaper = [box(s, 44, 200, 210, 120, fill=PAPER_BG, line=PAPER_BR, lw=1.2, radius=.04)]
gpaper.append(textbox(s, 58, 212, 180, 16, [[("drawn on the chart · p.10", {"size": 8, "bold": True, "color": RGBColor(0x8A,0x7F,0x63)})]]))
pbx = box(s, 60, 234, 178, 70, fill=WHITE, line=PAPER_BR, lw=1.1, radius=.05)
set_text(pbx, [[("Installation of tourniquet. (Position: one palm's width above the wound…)", {"size": 8.5, "color": PAPER_TX})]])
gpaper.append(pbx)
gpaper.append(textbox(s, 262, 240, 40, 40, [[("→", {"size": 20, "bold": True, "color": INK2})]]))
gcode = [codebox(s, 306, 180, 330, 208, [
    [('"hem_05": {', {})], [('  "type": "instruction",', {})],
    [('  "text_en": "Installation of', {})], [('              tourniquet. …",', {})],
    [('  "text_ar": "…تركيب الرباط الضاغط",', {})], [('  "page": 10,', {})],
    [('  "provenance": "source",', {"bold": True, "color": PRIMARY_D})],
    [('  "next": "hem_06"', {})], [('}', {})]], size=10)]
b1 = bullet(s, 690, 170, 550, [("Visual-only extraction", {"bold": True}),
    (" — pages rasterized @180 DPI and read as images. Text extraction drops the arrows, and ", {}),
    ("the arrows ARE the logic", {"bold": True}), (".", {})], size=14)
b2 = bullet(s, 690, 260, 550, [("Every node cites its ", {}), ("source page", {"bold": True}), (".", {})], size=14)
b3 = bullet(s, 690, 310, 550, [("Provenance audit trail: ", {}), ("58 of 59 nodes are source", {"bold": True}),
    ("; exactly 1 is added (navigation-only).", {})], size=14)
b4 = bullet(s, 690, 380, 550, [("validate.py ", {"font": MONO}), ("fails loudly", {"bold": True}),
    (": broken edge, orphan node, bad leaf → non-zero exit.", {})], size=14)
slide_no(s, 9)
animate(s, [gpaper, gcode, b1, b2, b3, b4])

# s10 — engine in Experta terms
s = prs.slides.add_slide(BLANK); SLIDES["s10"] = s
header(s, "The engine, in Experta terms", "Forward chaining = declare → match → fire")
code = [codebox(s, 44, 170, 620, 250, [
    [("# the baseline walker — engine.py (real code, elided)", {"color": INK3})],
    [("@Rule", {"bold": True, "color": PRIMARY_D}), ("(AS.p << ", {}), ("Position(nid=MATCH.nid)", {"bold": True, "color": PRIMARY_D}), (",", {})],
    [("      AS.a << ", {}), ("Active(chart=MATCH.c)", {"bold": True, "color": PRIMARY_D}), (",", {})],
    [("      NOT(Bleeding()), ", {}), ("salience=0", {"bold": True, "color": PRIMARY_D}), (")", {})],
    [("def", {"bold": True, "color": PRIMARY_D}), (" walk(self, p, a, nid, c):", {})],
    [("    node = self.nodes[nid]    # ← read from the JSON", {})],
    [("    …                         # render node, ask", {"color": INK3})],
    [("    self._move(p, a, next_id) # retract old Position,", {})],
    [("                # declare the new one → RETE re-matches", {"color": INK3})]], size=11)]
facts = []
fx = 44
for t in ["Position", "Active", "Danger", "Bleeding", "StillBleeding", "SignOfLife"]:
    w = 26 + len(t)*7.5
    facts.append(chip(s, fx, 440, w, 30, [(t, {"size": 11, "font": MONO, "color": INK2})], fill=SURFACE))
    fx += w + 10
facts.append(textbox(s, fx+6, 444, 320, 24, [[("— the six Fact types = the whole working memory", {"size": 10.5, "color": INK3})]]))
b1 = bullet(s, 700, 180, 540, [("One ", {}), ("KnowledgeEngine", {"font": MONO}), (" subclass; charts are ", {}),
    ("facts about position", {"bold": True}), (", not rules.", {})], size=14.5)
b2 = bullet(s, 700, 250, 540, [("Each answer ", {}), ("retracts / declares", {"bold": True}),
    (" a fact → the RETE network re-matches → the next rule fires. ", {}),
    ("That chain is the inference.", {"bold": True})], size=14.5)
b3 = bullet(s, 700, 345, 540, [("8 rules total", {"bold": True}), (" walk all 59 nodes.", {})], size=14.5)
tests = chip(s, 700, 400, 480, 36, [("11 automated tests — every inference behavior is proven, not claimed",
        {"size": 11.5, "bold": True, "color": SUCCESS})], fill=SUCCESS_T, line=SUCCESS)
slide_no(s, 10)
animate(s, [code, [*facts], b1, b2, b3, [tests]])

# s11 — inference I: salience
s = prs.slides.add_slide(BLANK); SLIDES["s11"] = s
header(s, "Inference I — salience & conflict resolution", "Salience encodes clinical precedence")
rungs = []
rdata = [("100", "ov_bleeding", "catastrophic bleeding — key b · bleed-out kills fastest", PRIMARY),
         ("90", "ov_arrest", "not breathing / no pulse — keys n · p", RGBColor(0xE0,0x55,0x55)),
         ("80", "ov_unconscious", "unconscious — key u", RGBColor(0xEF,0xA1,0xA1)),
         ("0", "walk", "the baseline chart walker", BORDER)]
ry = 170
for sal, nm, ds, bc in rdata:
    r = box(s, 44, ry, 560, 64, fill=WHITE, line=bc, lw=2, radius=.1)
    tf = r.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _style_runs(p, [(f" {sal}   ", {"size": 19, "bold": True, "color": PRIMARY if sal != "0" else INK3}),
                    (nm, {"size": 14, "bold": True, "font": MONO}),
                    (f"   {ds}", {"size": 10.5, "color": INK3})], 12)
    rungs.append([r]); ry += 76
b1 = bullet(s, 650, 165, 590, [("At ", {}), ("every", {"bold": True}), (" prompt, the paramedic can assert a ", {}),
    ("Danger", {"font": MONO}), (" fact ", {}), ("instead of answering", {"bold": True}), (".", {})], size=14)
b2 = bullet(s, 650, 230, 590, [("Two dangers at once? ", {}), ("Conflict resolution", {"bold": True}),
    (": the highest salience fires first — proven in TEST 1, not hand-sequenced.", {})], size=14)
b3 = bullet(s, 650, 315, 590, [("Guards: each override ", {}), ("retract", {"font": MONO}),
    ("s its Danger (fires once) and never re-enters the protocol it's already in.", {})], size=14)
tr = [codebox(s, 650, 400, 590, 78, [
    [("-- DANGER (type letter): [b]leeding  [n]ot-breathing  no-[p]ulse  [u]nconscious", {"color": INK2})],
    [("!! DANGER [catastrophic_bleeding] ", {"color": INK2}), ("(sal 100): OVERRIDE", {"bold": True}),
     (" -> preempt, jump to 'hemorrhage'", {"color": INK2})]], size=9),
    textbox(s, 650, 482, 500, 18, [[("verbatim engine output · presentation/traces/override_bleeding_en.txt",
        {"size": 8.5, "color": INK3})]])]
st = statement(s, 650, 520, 590, "A paper chart cannot accept an asynchronous fact.", size=17)
slide_no(s, 11)
animate(s, rungs + [b1, b2, b3, tr, st])

# s12 — inference II: CF
s = prs.slides.add_slide(BLANK); SLIDES["s12"] = s
header(s, "Inference II — certainty factor", "Uncertainty is data too")
chips_g = [chip(s, 44, 150, 240, 32, [("cpr_11 — “[1] Sign of life?”", {"size": 11.5, "bold": True, "color": INK2})]),
           chip(s, 300, 150, 500, 32, [("“Gasping is not considered efficient breathing” — chart #3, p.3",
                {"size": 11.5, "bold": True, "color": WARN_TX})], fill=WARN_T, line=WARN_TX)]
barg = []
barg.append(box(s, 60, 268, 869, 10, fill=PRIMARY_T, line=BORDER, lw=1, radius=.5))
barg.append(box(s, 929, 268, 291, 10, fill=SUCCESS_T, line=BORDER, lw=1, radius=.5))
barg.append(box(s, 927, 246, 3, 52, fill=INK, line=None, radius=0, shape=MSO_SHAPE.RECTANGLE))
barg.append(textbox(s, 795, 222, 125, 20, [[("threshold 0.5", {"size": 11, "bold": True})]], align=PP_ALIGN.RIGHT))
for label, cf, frac, col in [("none", "CF −1.0", .02, PRIMARY), ("unsure", "CF +0.2", .59, PRIMARY),
                             ("likely", "CF +0.6", .795, SUCCESS), ("certain", "CF +1.0", .98, SUCCESS)]:
    mx = 60 + frac*1160
    dot = box(s, mx-7, 266, 14, 14, fill=col, line=WHITE, lw=2, shape=MSO_SHAPE.OVAL, radius=None)
    barg += [dot, textbox(s, mx-55, 226, 110, 36, [[(label, {"size": 12, "bold": True})], [(cf, {"size": 9.5, "color": INK3})]],
             align=PP_ALIGN.CENTER)]
branches = [textbox(s, 60, 300, 620, 24, [[("← below 0.5 → CPR branch (cpr_14) — “if in doubt, compress”",
             {"size": 12.5, "bold": True, "color": PRIMARY_D})]]),
            textbox(s, 700, 300, 520, 24, [[("≥ 0.5 → ventilation branch (cpr_12) →",
             {"size": 12.5, "bold": True, "color": SUCCESS})]], align=PP_ALIGN.RIGHT)]
b1 = bullet(s, 44, 370, 600, [("Experta has ", {}), ("no native CF", {"bold": True}), (" → we model it as data: a ", {}),
    ("SignOfLife(cf)", {"font": MONO}), (" fact + ", {}), ("one thresholding rule", {"bold": True}), (" (salience 20).", {})], size=14)
b2 = bullet(s, 44, 455, 600, [("The conservative bias lives ", {}), ("in the scale", {"bold": True}),
    (": genuine doubt scores 0.2, far below the cut. Endpoints match a plain yes/no — only the ambiguous middle re-routes.", {})], size=14)
tr = [codebox(s, 690, 380, 550, 92, [
    [("?? Confidence the patient has a SIGN OF LIFE?", {"color": INK2})],
    [("3) unsure  (CF +0.20)", {"color": INK2})],
    [("~ CF +0.20 <  +0.50: ", {"color": INK2}), ("doubt -> CPR branch", {"bold": True}),
     (" (if in doubt, compress)", {"color": INK2})]], size=9.5),
    textbox(s, 690, 476, 480, 18, [[("verbatim engine output · presentation/traces/cf_unsure_en.txt", {"size": 8.5, "color": INK3})]])]
slide_no(s, 12)
animate(s, [chips_g, barg, branches, b1, b2, tr])

# ── walkthrough tree (ported 1:1 from deck.html) ───────────────────────────
TX, TY = 54, 118   # tree origin on the pptx slide
NODES = {
 "gen01": (0, 26, 128, 44, "", "Safety", "gen_01"),
 "gen02": (146, 26, 132, 44, "", "Scene", "gen_02"),
 "gen03": (296, 26, 132, 44, "", "Situation", "gen_03"),
 "gen04": (446, 20, 196, 56, "q", "◇ Life-threatening bleeding?", "gen_04"),
 "gen05": (662, 24, 148, 52, "jump", "⇢ Hemorrhage", "gen_05 · jump"),
 "hem01": (40, 130, 292, 64, "", "Direct pressure · compression dressing", "hem_01"),
 "hem02": (370, 132, 196, 60, "q", "◇ Active bleeding?", "hem_02"),
 "hem03": (40, 250, 292, 64, "", "Second compression bandage", "hem_03 · TIER 2"),
 "hem04": (370, 252, 196, 60, "q", "◇ Active bleeding?", "hem_04"),
 "hem05": (40, 370, 292, 64, "", "Tourniquet — one palm above wound", "hem_05 · terminal"),
 "hem06": (370, 372, 196, 60, "q", "◇ Active bleeding?", "hem_06"),
 "hem08": (616, 186, 224, 50, "leaf", "Transport to health center", "hem_08"),
 "hem07": (616, 378, 224, 50, "leaf", "URGENT transport", "hem_07"),
}
EDGES = {
 "g12": "M128,47 L146,47", "g23": "M278,47 L296,47", "g34": "M428,47 L442,47",
 "g45": "M646,47 L657,47", "jmp": "M736,78 V88 H185 V124",
 "h12": "M332,162 L366,162", "y1": "M466,196 V225 H185 V244",
 "n1": "M566,162 H596 V206 H616", "h34": "M332,282 L366,282",
 "y2": "M466,316 V345 H185 V364", "n2": "M566,282 H596 V236 H616",
 "h56": "M332,402 L366,402", "n3": "M566,402 L614,402", "hold": "M466,436 V465 H700 V428",
}
LABELS = {
 "ljmp": (310, 76, 330, 24, "yes →  >> JUMP to chart 'hemorrhage' <<", INK2, WHITE),
 "ly1":  (250, 212, 170, 24, "yes — still bleeding", PRIMARY, WHITE),
 "ly2":  (250, 332, 170, 24, "yes — still bleeding", PRIMARY, WHITE),
 "ln1":  (596, 118, 150, 24, "✓ no — controlled", SUCCESS, WHITE),
 "ln2":  (596, 298, 150, 24, "✓ no — controlled", SUCCESS, WHITE),
 "ln3":  (572, 426, 66, 22, "✓ no", SUCCESS, WHITE),
 "lhold":(196, 452, 430, 24, "⚠ yes at terminal tier → HOLD, no tier 4 (flagged divergence)", WARN_TX, WARN_T),
}
EDGE_COLOR = {"future": (GRAY_EDGE, 1.5), "visited": (INK2, 2.0),
              "delta": (PRIMARY, 2.5), "good": (SUCCESS, 2.5)}

def draw_tree(s, nstates, estates, labstates, tiers):
    textbox(s, TX, TY, 400, 16, [[("GENERAL ASSESSMENT — CHART #1", {"size": 8.5, "bold": True, "color": INK3})]])
    textbox(s, TX, TY+104, 440, 16, [[("CATASTROPHIC HEMORRHAGE — CHART #10", {"size": 8.5, "bold": True, "color": INK3})]])
    for eid, d in EDGES.items():
        col, wd = EDGE_COLOR[estates.get(eid, "future")]
        dd = re.sub(r"([\d.]+)", lambda m: str(float(m.group(1))), d)
        for i, (x1, y1, x2, y2) in enumerate(path_segs(d)):
            line_seg(s, TX+x1, TY+y1, TX+x2, TY+y2, col, wd, arrow=(i == len(path_segs(d))-1))
    for nid_, (x, y, w, h, kind, label, sub) in NODES.items():
        st = nstates.get(nid_, "future")
        fill, lc, lw_, tcol, bold = WHITE, BORDER, 1.5, INK3, False
        if st == "visited": lc, tcol, lw_, bold = INK2, INK, 2.0, True
        if st == "current": fill, lc, lw_, tcol, bold = PRIMARY_T, PRIMARY, 2.5, INK, True
        if st == "current-good": fill, lc, lw_, tcol, bold = SUCCESS_T, SUCCESS, 2.5, INK, True
        shp = box(s, TX+x, TY+y, w, h, fill=fill, line=lc, lw=lw_,
                  radius=(.5 if kind == "leaf" else .14), dash=("dash" if kind == "jump" else None))
        set_text(shp, [[(label, {"size": 11, "bold": bold, "color": tcol})],
                       [(sub, {"size": 7.5, "color": INK3})]])
    for i, on in enumerate(tiers):
        bshp = box(s, TX, TY+148+i*120, 28, 28, fill=PRIMARY if on else SURFACE,
                   line=PRIMARY if on else BORDER, lw=1.5, shape=MSO_SHAPE.OVAL, radius=None)
        set_text(bshp, [[(str(i+1), {"size": 11, "bold": True, "color": WHITE if on else INK2})]])
    for lid, (x, y, w, h, txt, col, fl) in LABELS.items():
        st = labstates.get(lid, "dim")
        if st == "off": continue
        c = col if st == "on" else {PRIMARY: DIM_RED, SUCCESS: DIM_GREEN, WARN_TX: DIM_WARN, INK2: GRAY_EDGE}.get(col, col)
        lb = box(s, TX+x, TY+y, w, h, fill=fl if st == "on" else WHITE, line=c, lw=1.25, radius=.4)
        set_text(lb, [[(txt, {"size": 9, "bold": True, "color": c})]])

def side_panel(s, facts=None, rule=None, deltas=None, note=None, out="", src="", vote=False):
    box(s, 932, TY, 304, 552, fill=WHITE, line=BORDER, lw=1.5, radius=.05)
    textbox(s, 950, TY+14, 260, 20, [[("ENGINE STATE", {"size": 10, "bold": True})]])
    line_seg(s, 950, TY+38, 1218, TY+38, BORDER, 1)
    y = TY + 48
    textbox(s, 950, y, 260, 16, [[("WORKING MEMORY", {"size": 8, "bold": True, "color": INK3})]]); y += 20
    if vote:
        for t in ["Position(nid=gen_04)", "Active(chart=general_…)"]:
            c = chip(s, 950, y, 250, 26, [(t, {"size": 10, "font": MONO, "color": INK2})], fill=SURFACE); y += 32
        y += 8
        textbox(s, 950, y, 260, 16, [[("YOUR ANSWER BECOMES A FACT", {"size": 8, "bold": True, "color": INK3})]]); y += 22
        yes = box(s, 950, y, 268, 84, fill=PRIMARY, line=None, radius=.12)
        set_text(yes, [[("YES", {"size": 22, "bold": True, "color": WHITE})],
                       [("blood is spreading fast", {"size": 10.5, "bold": True, "color": PRIMARY_T})]])
        LINKS.append((yes, "w5")); y += 96
        no = box(s, 950, y, 268, 84, fill=WHITE, line=SUCCESS, lw=2, radius=.12)
        set_text(no, [[("NO", {"size": 22, "bold": True, "color": SUCCESS})],
                      [("no life-threatening bleed", {"size": 10.5, "bold": True, "color": SUCCESS})]])
        LINKS.append((no, "wno")); y += 100
    else:
        for t, hot in facts:
            c = chip(s, 950, y, min(258, 30+len(t)*6.4), 26,
                     [(t, {"size": 10, "font": MONO, "bold": bool(hot), "color": PRIMARY_D if hot else INK2})],
                     fill=PRIMARY_T if hot else SURFACE, line=PRIMARY if hot else BORDER)
            y += 32
        y += 6
        textbox(s, 950, y, 260, 16, [[("RULE FIRED", {"size": 8, "bold": True, "color": INK3})]]); y += 18
        textbox(s, 950, y, 268, 24, [[(rule[0], {"size": 13, "bold": True, "font": MONO}),
                                      ("   " + rule[1], {"size": 9.5, "color": INK2})]]); y += 32
        textbox(s, 950, y, 260, 16, [[("Δ THIS STEP", {"size": 8, "bold": True, "color": INK3})]]); y += 18
        for op, t, nt in deltas:
            neg = op == "-"
            textbox(s, 950, y, 268, 20, [[
                ("− " if neg else "+ ", {"size": 11, "bold": True, "color": INK3 if neg else PRIMARY, "font": MONO}),
                (t, {"size": 11, "font": MONO, "strike": neg, "bold": not neg,
                     "color": INK3 if neg else PRIMARY_D}),
                ("  " + nt, {"size": 8.5, "color": INK3})]])
            y += 24
        if note:
            y += 6
            box(s, 950, y, 4, 44, fill=PRIMARY, line=None, radius=.3)
            textbox(s, 962, y, 256, 52, [[(note, {"size": 10, "bold": True})]], spacing=1.1); y += 58
    oy = TY + 552 - 118
    textbox(s, 950, oy-18, 260, 16, [[("ENGINE OUTPUT (VERBATIM)", {"size": 8, "bold": True, "color": INK3})]])
    codebox(s, 950, oy, 268, 84, [[ (ln, {"color": INK2}) ] for ln in out.split("\n")], size=8.5)
    textbox(s, 950, oy+88, 268, 24, [[("real session · presentation/traces/" + src, {"size": 7.5, "color": INK3})]])

def walk_slide(no, key, kicker, title, nstates, estates, labstates, tiers, dots,
               facts=None, rule=None, deltas=None, note=None, out="", src="", vote=False):
    s = prs.slides.add_slide(BLANK); SLIDES[key] = s
    header(s, kicker, title, tsize=21)
    draw_tree(s, nstates, estates, labstates, tiers)
    side_panel(s, facts, rule, deltas, note, out, src, vote)
    chip(s, 44, 682, 520, 26, [("SARC charts #1 + #10 · data/general_assessment.json · data/hemorrhage.json",
         {"size": 9, "color": INK2})], fill=SURFACE)
    dx = 1050
    for i in range(10):
        col = PRIMARY if i+1 == dots else (INK2 if i+1 < dots else BORDER)
        box(s, dx, 690, 10 if i+1 == dots else 8, 10 if i+1 == dots else 8,
            fill=col, line=None, shape=MSO_SHAPE.OVAL, radius=None)
        dx += 17
    slide_no(s, no)
    return s

V, C, F = "visited", "current", "future"
walk_slide(13, "w1", "Flow walkthrough · step 1 of 10", "Two facts. One generic rule.",
    {"gen01": C}, {}, {}, [0,0,0], 1,
    facts=[("Position(nid=gen_01)",1),("Active(chart=general_…)",1)], rule=("walk","salience 0"),
    deltas=[("+","Position(nid=gen_01)","declare"),("+","Active(chart=…)","declare")],
    out="[gen_01] Safety: ensure safety &\nsecurity for staff, patient(s)\nand bystanders…  (p.1)", src="walkthrough_yes_branch_en.txt")
walk_slide(14, "w2", "Flow walkthrough · step 2 of 10", "Answer = retract, declare. RETE re-matches.",
    {"gen01": V, "gen02": C}, {"g12": "delta"}, {}, [0,0,0], 2,
    facts=[("Position(nid=gen_02)",1),("Active(chart=general_…)",0)], rule=("walk","salience 0"),
    deltas=[("-","Position(nid=gen_01)","retract"),("+","Position(nid=gen_02)","declare")],
    out="[gen_02] Scene Evaluation:\nmechanism of injury;\nenvironmental factors…  (p.1)", src="walkthrough_yes_branch_en.txt")
walk_slide(15, "w3", "Flow walkthrough · step 3 of 10", "That was the Preliminary phase — SSS.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": C}, {"g12": V, "g23": V, "g34": "delta"}, {}, [0,0,0], 3,
    facts=[("Position(nid=gen_04)",1),("Active(chart=general_…)",0)], rule=("walk","salience 0 · fired twice"),
    deltas=[("-","Position(nid=gen_03)","retract"),("+","Position(nid=gen_04)","declare")],
    out="[gen_03] Situation Evaluation…\n[gen_04] Life threatening\nbleeding?   1) yes  2) no", src="walkthrough_yes_branch_en.txt")
walk_slide(16, "w4", "Flow walkthrough · step 4 of 10", "The chart's most important question. You decide.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": C}, {"g12": V, "g23": V, "g34": V}, {}, [0,0,0], 4,
    out="[gen_04] Life threatening\nbleeding?\n   1) yes   2) no\n   > _", src="walkthrough_yes_branch_en.txt", vote=True)
walk_slide(17, "w5", "Flow walkthrough · step 5 of 10", "A jump: the knowledge is five separate trees.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": V, "gen05": V, "hem01": C},
    {"g12": V, "g23": V, "g34": V, "g45": V, "jmp": "delta"}, {"ljmp": "on"}, [0,0,0], 5,
    facts=[("Position(nid=hem_01)",1),("Active(chart=hemorrhage)",1)], rule=("walk","salience 0 · resolves the jump"),
    deltas=[("-","Active(chart=general_…)","retract"),("+","Active(chart=hemorrhage)","declare")],
    note="ABCDE never ran — circulation-first, as structure.",
    out="[gen_05] Catastrophic Hemorrhage.\n>> JUMP to chart 'hemorrhage'\n(node 'hem_01') <<", src="walkthrough_yes_branch_en.txt")
walk_slide(18, "w6", "Flow walkthrough · step 6 of 10", "A NEW kind of fact: state, not position.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": V, "gen05": V, "hem01": C},
    {"g12": V, "g23": V, "g34": V, "g45": V, "jmp": V}, {"ljmp": "on"}, [1,0,0], 6,
    facts=[("Position(nid=hem_01)",0),("Active(chart=hemorrhage)",0),("Bleeding(tier=1)",1)],
    rule=("hem_enter","salience 50"), deltas=[("+","Bleeding(tier=1)","declare")],
    note="The escalation tier now lives in working memory.",
    out=">> entering hemorrhage escalation\n(semantic tier loop) <<\n>> ESCALATION TIER 1: Control of\nexternal bleeding…", src="walkthrough_yes_branch_en.txt")
walk_slide(19, "w7", "Flow walkthrough · step 7 of 10", "Escalation: the fact remembers.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": V, "gen05": V, "hem01": V, "hem02": V, "hem03": C},
    {"g12": V, "g23": V, "g34": V, "g45": V, "jmp": V, "h12": V, "y1": "delta"},
    {"ljmp": "on", "ly1": "on"}, [1,1,0], 7,
    facts=[("Position(nid=hem_03)",0),("Active(chart=hemorrhage)",0),("Bleeding(tier=2)",1)],
    rule=("hem_decide","salience 10"),
    deltas=[("-","Bleeding(tier=1)","retract"),("+","Bleeding(tier=2)","declare")],
    note="It never re-asks tier 1 — the fact remembers.",
    out=">> still bleeding -> ADVANCE\ntier 1 -> 2\n>> ESCALATION TIER 2: Second\ncompression bandage…", src="walkthrough_yes_branch_en.txt")
walk_slide(20, "w8", "Flow walkthrough · step 8 of 10", "Tier 3 is terminal: the engine refuses to spin.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": V, "gen05": V, "hem01": V, "hem02": V, "hem03": V, "hem04": V, "hem05": C},
    {"g12": V, "g23": V, "g34": V, "g45": V, "jmp": V, "h12": V, "y1": V, "h34": V, "y2": "delta"},
    {"ljmp": "on", "ly1": "on", "ly2": "on", "lhold": "on"}, [1,1,1], 8,
    facts=[("Position(nid=hem_05)",0),("Active(chart=hemorrhage)",0),("Bleeding(tier=3)",1)],
    rule=("hem_decide","salience 10"),
    deltas=[("-","Bleeding(tier=2)","retract"),("+","Bleeding(tier=3)","declare")],
    note="Chart draws an endless re-apply loop; we cap at terminal — a flagged divergence.",
    out=">> still bleeding -> ADVANCE\ntier 2 -> 3\n>> ESCALATION TIER 3:\nInstallation of tourniquet…", src="walkthrough_yes_branch_en.txt")
walk_slide(21, "w9", "Flow walkthrough · step 9 of 10", "Controlled → leaf. Zero chart-specific code.",
    {"gen01": V, "gen02": V, "gen03": V, "gen04": V, "gen05": V, "hem01": V, "hem02": V, "hem03": V,
     "hem04": V, "hem05": V, "hem06": V, "hem07": "current-good"},
    {"g12": V, "g23": V, "g34": V, "g45": V, "jmp": V, "h12": V, "y1": V, "h34": V, "y2": V, "h56": V, "n3": "good"},
    {"ljmp": "on", "ly1": "on", "ly2": "on", "ln3": "on"}, [1,1,1], 9,
    facts=[("Position(nid=hem_07)",1),("Active(chart=hemorrhage)",0)],
    rule=("hem_decide → walk","exit, then leaf"),
    deltas=[("-","Bleeding(tier=3)","retract"),("-","StillBleeding(no)","retract")],
    out=">> bleeding controlled at tier 3\n-> exit to hem_07\n[hem_07] Urgent and safe\ntransport…  === END (leaf) ===", src="walkthrough_yes_branch_en.txt")

# s19 (key) — the NO branch strip
s = prs.slides.add_slide(BLANK); SLIDES["wno"] = s
header(s, "Flow walkthrough · step 10 of 10", "The road not taken — the NO branch", tsize=21)
strip = [("◇ Life-threatening bleeding?", "gen_04", "q"), ("✓ no", "", "lab"),
         ("Primary Assessment — AVPU + ABCDE", "gen_06", ""), ("Stable or unstable?", "gen_07", "q"),
         ("Unstable: immediate action, BLS", "gen_09", ""),
         ("Secondary Assessment — SAMPLER · OPQRST", "gen_10", ""),
         ("Continuous Assessment & Orientation", "gen_11", "leafg")]
x = 44
for label, sub, kind in strip:
    if kind == "lab":
        chip(s, x, 208, 74, 28, [(label, {"size": 10, "bold": True, "color": SUCCESS})], line=SUCCESS)
        x += 90; continue
    w = 168 if kind != "leafg" else 190
    fill = SUCCESS_T if kind == "leafg" else WHITE
    lc = SUCCESS if kind == "leafg" else INK2
    bx = box(s, x, 180, w, 84, fill=fill, line=lc, lw=2, radius=(.5 if kind == "leafg" else .12))
    set_text(bx, [[(label, {"size": 10.5, "bold": True})], [(sub, {"size": 8, "color": INK3})]])
    x += w + 8
    if kind != "leafg" and label != strip[-1][0]:
        textbox(s, x-4, 208, 20, 24, [[("→", {"size": 13, "bold": True, "color": INK3})]]); x += 20
codebox(s, 44, 330, 640, 128, [
    [("[gen_06] Primary Assessment — Look-Listen-Feel-Do. AVPU… ABCDE…", {"color": INK2})],
    [("[gen_07] Intervention Strategy — stable or unstable?", {"color": INK2})],
    [("[gen_10] Secondary Assessment — SAMPLER… OPQRST… GCS… PERRLA…", {"color": INK2})],
    [("[gen_11] Continuous Assessment & Orientation…", {"color": INK2})],
    [("=== END (action / leaf) ===", {"bold": True})]], size=9.5)
textbox(s, 44, 462, 500, 18, [[("real session · presentation/traces/walkthrough_no_branch_en.txt", {"size": 8.5, "color": INK3})]])
chip(s, 720, 340, 500, 36, [("No new fact types here — this is the plain walk at salience 0.", {"size": 11.5, "color": INK2})])
back = chip(s, 720, 390, 500, 36, [("▶ …and the YES branch is where the inference lives — watch it",
       {"size": 11.5, "bold": True, "color": PRIMARY_D})], fill=PRIMARY_T, line=PRIMARY)
LINKS.append((back, "w5"))
slide_no(s, 22)

# s24 — FastAPI
s = prs.slides.add_slide(BLANK); SLIDES["s24"] = s
header(s, "Integration — the API", "The API contract is a consequence of the architecture")
code = [codebox(s, 44, 170, 620, 300, [
    [("POST /sessions/{id}/step      ← real response (abridged)", {"color": INK3})],
    [("{", {})], [('  "screen": {', {})],
    [('    "type": "question",', {"bold": True, "color": PRIMARY_D}), ("        ← the whole contract", {"color": INK3})],
    [('    "node_id": "chk_01",', {})], [('    "chart_id": "choking",', {})],
    [('    "text": { "en": "Suspect foreign body…",', {})],
    [('              "ar": "…هل يُشتبه في انسداد" },', {})],
    [('    "options": [ {"id":1,"label":{"en":"yes","ar":"نعم"}},', {})],
    [('                 {"id":2,"label":{"en":"no","ar":"لا"}} ]', {})],
    [("  },", {})], [('  "danger_panel": [ …b · n · p · u… ]', {})], [("}", {})]], size=10)]
b1 = bullet(s, 700, 190, 540, [("Response shapes are driven by ", {}), ("node type from the JSON knowledge", {"bold": True}),
    (" — never by internal rules or salience.", {})], size=14.5)
b2 = bullet(s, 700, 275, 540, [("New rules → ", {}), ("zero API change", {"bold": True}), (". New chart → ", {}),
    ("auto-discovered", {"bold": True}), (" by GET /charts.", {})], size=14.5)
b3 = bullet(s, 700, 360, 540, [("Every text field ships ", {}), ("both languages", {"bold": True}),
    (" → switching EN↔AR needs no server round-trip.", {})], size=14.5)
slide_no(s, 23)
animate(s, [code, b1, b2, b3])

# s25 — Flutter app (three phone mockups)
s = prs.slides.add_slide(BLANK); SLIDES["s25"] = s
header(s, "Integration — the app", "The same knowledge, in a responder's pocket")
def phone(x, caption, body):
    g = []
    g.append(box(s, x, 160, 200, 400, fill=WHITE, line=INK, lw=2.5, radius=.12))
    hd = box(s, x+6, 166, 188, 34, fill=PRIMARY, line=None, radius=.15)
    set_text(hd, [[("ResQRules · " + caption[0], {"size": 8.5, "bold": True, "color": WHITE})]])
    g.append(hd); g += body(x)
    dp = box(s, x+6, 522, 188, 32, fill=PRIMARY, line=None, radius=.1)
    set_text(dp, [[("BLEEDING | NO BREATH | NO PULSE | UNCONSC.", {"size": 6.5, "bold": True, "color": WHITE})]])
    g.append(dp)
    g.append(textbox(s, x-10, 568, 220, 40, [[(caption[1], {"size": 11, "color": INK2})]], align=PP_ALIGN.CENTER))
    return g
def body1(x):
    q = textbox(s, x+14, 212, 172, 66, [[("Suspect foreign body obstruction? (Universal sign: hands to throat…)",
        {"size": 8.5, "bold": True})]], spacing=1.15)
    o1 = box(s, x+14, 286, 172, 30, fill=WHITE, line=BORDER, lw=1.25, radius=.25)
    set_text(o1, [[("yes", {"size": 9, "bold": True})]])
    o2 = box(s, x+14, 322, 172, 30, fill=WHITE, line=BORDER, lw=1.25, radius=.25)
    set_text(o2, [[("no", {"size": 9, "bold": True})]])
    cite = textbox(s, x+14, 494, 160, 16, [[("chk_01 · p.7", {"size": 7, "color": INK3})]])
    return [q, o1, o2, cite]
def body2(x):
    t = box(s, x+14, 212, 172, 44, fill=PRIMARY, line=None, radius=.18)
    set_text(t, [[("⚡ OVERRIDE → Catastrophic Hemorrhage · salience 100", {"size": 7.5, "bold": True, "color": WHITE})]])
    q = textbox(s, x+14, 266, 172, 60, [[("Control of external bleeding: direct compression on the wound…",
        {"size": 8.5, "bold": True})]], spacing=1.15)
    o = box(s, x+14, 330, 172, 30, fill=PRIMARY_T, line=PRIMARY, lw=1.25, radius=.25)
    set_text(o, [[("continue", {"size": 9, "bold": True})]])
    cite = textbox(s, x+14, 494, 172, 16, [[("hem_01 · p.10 · tier 1/3", {"size": 7, "color": INK3})]])
    return [t, q, o, cite]
def body3(x):
    lf = box(s, x+14, 212, 172, 66, fill=SUCCESS_T, line=SUCCESS, lw=1.5, radius=.18)
    set_text(lf, [[("✓ Urgent and safe transport to a health center.", {"size": 9, "bold": True, "color": SUCCESS})]])
    rs = box(s, x+14, 470, 172, 30, fill=WHITE, line=BORDER, lw=1.25, radius=.25)
    set_text(rs, [[("↺ new session", {"size": 9, "bold": True})]])
    return [lf, rs]
p1 = phone(60, ("Heimlich Maneuver", "question screen"), body1)
p2 = phone(300, ("Catastrophic Hemorrhage", "danger-panel override (Inference I, as UI)"), body2)
p3 = phone(540, ("Catastrophic Hemorrhage", "terminal action leaf"), body3)
b1 = bullet(s, 800, 200, 440, [("One widget per server node ", {}), ("type", {"font": MONO}),
    (" — ", {}), ("the KBS drives the UI", {"bold": True}), ("; the shell never changes.", {})], size=13.5)
b2 = bullet(s, 800, 280, 440, [("The persistent ", {}), ("danger panel", {"bold": True}),
    (" = the salience overrides, as UI.", {})], size=13.5)
b3 = bullet(s, 800, 345, 440, [("Instant EN↔AR — both languages arrive in every response; RTL flips locally.", {})], size=13.5)
b4 = bullet(s, 800, 420, 440, [("Stack, in one breath: Clean Architecture · Cubit · Retrofit.", {"size": 11.5, "color": INK2})], size=11.5)
slide_no(s, 24)
animate(s, [p1, p2, p3, b1, b2, b3, b4])

# s26 — demo (red)
s = red_slide(); SLIDES["s26"] = s
textbox(s, 190, 210, 900, 90, [[("Live demo", {"size": 54, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
textbox(s, 240, 320, 800, 40, [[("The engine narrates its own inference — watch the salience fire.",
    {"size": 17, "color": RGBColor(0xFF, 0xE2, 0xE2)})]], align=PP_ALIGN.CENTER)
cxx = 170
for t in ["danger key b → sal-100 preempt", "“unsure” at Sign of life? → CF routes to compressions", "the app, in the field"]:
    w = 40 + len(t)*6.8
    cbox = box(s, cxx, 400, w, 36, fill=RGBColor(0xD9, 0x33, 0x33), line=RGBColor(0xF0, 0xA0, 0xA0), lw=1, radius=.5)
    set_text(cbox, [[(t, {"size": 10.5, "bold": True, "color": WHITE})]])
    cxx += w + 20
textbox(s, 240, 480, 800, 30, [[("(presenter: switch to the video / live terminal now)",
    {"size": 12, "color": RGBColor(0xFF, 0xB9, 0xB9)})]], align=PP_ALIGN.CENTER)
slide_no(s, 25)

# s27 — recap
s = prs.slides.add_slide(BLANK); SLIDES["s27"] = s
header(s, "The journey", "From paper to a reasoning system")
steps = [("SARC PDFs", "visual-only extraction"), ("JSON knowledge", "59 nodes · provenance"),
         ("Experta engine", "8 rules · 3 layers"), ("FastAPI", "type-driven contract"),
         ("Flutter app", "bilingual · danger panel")]
g_steps = []
x = 44
for i, (t, d) in enumerate(steps):
    bx = box(s, x, 210, 212, 86, fill=WHITE, line=INK2, lw=1.75, radius=.12)
    set_text(bx, [[(t, {"size": 13.5, "bold": True})], [(d, {"size": 9.5, "color": INK3})]])
    grp = [bx]
    if i < 4: grp.append(textbox(s, x+214, 236, 26, 30, [[("→", {"size": 15, "bold": True, "color": INK3})]]))
    g_steps.append(grp); x += 246
thesis = textbox(s, 140, 360, 1000, 110, [
    [("Paper can't react.", {"size": 26, "bold": True, "color": PRIMARY_D})],
    [("Knowledge as data + a small inference engine can.", {"size": 26, "bold": True, "color": PRIMARY_D})]],
    align=PP_ALIGN.CENTER, spacing=1.2)
limits = chip(s, 250, 510, 780, 36, [("honest limits: jumps are one-way (no return-stack yet) · Phase-2 charts are JSON-only additions",
    {"size": 11.5, "color": INK2})], fill=SURFACE)
slide_no(s, 26)
animate(s, g_steps + [[thesis], [limits]])

# s28 — team
s = prs.slides.add_slide(BLANK); SLIDES["s28"] = s
header(s, "The team — in journey order", "Knowledge → engine → API → app")
team = [("M.GH", "KNOWLEDGE", "Software engineer with paramedic experience · knowledge source & reference collector · Flutter dev"),
        ("M.SH", "ENGINE", "Python expert · AI engineer — the Experta inference engine"),
        ("M.JA", "API", "Backend engineer — wrapped the engine in FastAPI"),
        ("M.Y.D", "APP", "Flutter developer & designer"),
        ("M.Marwan", "APP · OPS", "Flutter expert · UI designer · DevOps on the server · Linux lover")]
g_cards = []
x = 44
for nm, tag, ds in team:
    card = box(s, x, 190, 228, 260, fill=WHITE, line=BORDER, lw=1.75, radius=.08)
    av = box(s, x+84, 210, 60, 60, fill=PRIMARY_T, line=PRIMARY, lw=2, shape=MSO_SHAPE.OVAL, radius=None)
    set_text(av, [[(nm.replace("M.Marwan", "M.M"), {"size": 11, "bold": True, "color": PRIMARY_D})]])
    tg = textbox(s, x+14, 282, 200, 18, [[(tag, {"size": 9, "bold": True, "color": PRIMARY_D})]], align=PP_ALIGN.CENTER)
    nn = textbox(s, x+14, 302, 200, 24, [[(nm, {"size": 14, "bold": True})]], align=PP_ALIGN.CENTER)
    dd = textbox(s, x+14, 330, 200, 110, [[(ds, {"size": 10, "color": INK2})]], align=PP_ALIGN.CENTER, spacing=1.2)
    g_cards.append([card, av, tg, nn, dd]); x += 240
sup = chip(s, 44, 660, 260, 30, [("Supervised by Miss Omama", {"size": 11, "color": INK2})], fill=SURFACE)
slide_no(s, 27)
animate(s, g_cards + [[sup]])

# s29 — thanks (red)
s = red_slide(); SLIDES["s29"] = s
textbox(s, 190, 220, 900, 60, [[("Thank you, Miss Omama.", {"size": 34, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
textbox(s, 240, 300, 800, 34, [[("For the guidance that kept the knowledge honest and the inference real.",
    {"size": 15, "color": RGBColor(0xFF, 0xE2, 0xE2)})]], align=PP_ALIGN.CENTER)
textbox(s, 190, 380, 900, 80, [[("Questions?", {"size": 46, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER)
textbox(s, 240, 480, 800, 26, [[("ResQRules · 8 rules · 59 nodes · 5 protocols",
    {"size": 12, "color": RGBColor(0xFF, 0xB9, 0xB9)})]], align=PP_ALIGN.CENTER)
slide_no(s, 28)

# renumber: we created 28 slides (walkthrough C1 shares numbering); fix page labels is cosmetic.

# resolve hyperlinks
for shape, key in LINKS:
    shape.click_action.target_slide = SLIDES[key]

OUT = pathlib.Path(__file__).parent / "deck.pptx"
prs.save(str(OUT))
print("slides:", len(prs.slides._sldIdLst))
print("saved", OUT)
